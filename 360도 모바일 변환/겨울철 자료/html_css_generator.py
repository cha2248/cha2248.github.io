# 좋아 마음에 든다. 테스트 완료!
# 아, 이미지 경로에 대해 물어볼 필요는 없을꺼같아. 그부분은 지워줘
# 자동화 완료!







import os
from pptx import Presentation
from pptx.enum.action import PP_ACTION

# 객체 위치 계산 함수
def calculate_position(slide_width, slide_height, shape):
    center_x = shape.left + (shape.width / 2)
    center_y = shape.top + (shape.height / 2)
    pos_x_percent = (center_x / slide_width) * 100
    pos_y_percent = (center_y / slide_height) * 100
    return round(pos_x_percent, 2), round(pos_y_percent, 2)

# 개별 객체에서 하이퍼링크 추출 함수
def extract_from_shape(shape, slide_num, slide_width, slide_height, hyperlinks):
    if shape.shape_type == 6:  # 그룹 객체
        for sub_shape in shape.shapes:  # 하위 객체 순회
            extract_from_shape(sub_shape, slide_num, slide_width, slide_height, hyperlinks)
    else:
        # 일반 객체의 하이퍼링크 확인
        if hasattr(shape, 'click_action') and shape.click_action and shape.click_action.action == PP_ACTION.HYPERLINK:
            address = shape.click_action.hyperlink.address
            if address and address.startswith(('http://', 'https://')):
                pos_x, pos_y = calculate_position(slide_width, slide_height, shape)
                hyperlinks.append((slide_num, pos_x, pos_y, address))
        # 텍스트 프레임 내 하이퍼링크 확인
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    address = run.hyperlink.address
                    if address and address.startswith(('http://', 'https://')):
                        pos_x, pos_y = calculate_position(slide_width, slide_height, shape)
                        hyperlinks.append((slide_num, pos_x, pos_y, address))

# 하이퍼링크 추출 함수
def extract_hyperlinks(pptx_file_path):
    hyperlinks = []
    prs = Presentation(pptx_file_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            extract_from_shape(shape, slide_num, slide_width, slide_height, hyperlinks)
    
    return hyperlinks

# 1. 사용자로부터 입력받기
ppt_path = input("파워포인트 파일 경로를 입력해주세요: ")
title_name = input("타이틀 명칭을 입력해주세요 (예: 함평군공립요양병원): ")

# 2. 파일명 및 경로 추출
try:
    ppt_dir = os.path.dirname(ppt_path)  # 파워포인트 파일이 있는 디렉토리
    file_name = os.path.splitext(os.path.basename(ppt_path))[0]  # 확장자 제외 파일명
    output_html_path = os.path.join(ppt_dir, f"{file_name}.html")  # 출력 HTML 파일 경로
    output_css_path = os.path.join(ppt_dir, f"{file_name}icons.css")  # 출력 CSS 파일 경로
    print(f"추출된 파일명: {file_name}")
    print(f"HTML 파일 저장 경로: {output_html_path}")
    print(f"CSS 파일 저장 경로: {output_css_path}")
except Exception as e:
    print(f"파일 경로 처리 중 오류 발생: {e}")
    exit()

# 3. 파워포인트 파일 열고 슬라이드 수 확인
try:
    prs = Presentation(ppt_path)
    slide_count = len(prs.slides)
    print(f"슬라이드 개수: {slide_count}")
except Exception as e:
    print(f"파워포인트 파일을 열 수 없습니다: {e}")
    exit()

# 4. 하이퍼링크 정보 추출
hyperlinks = extract_hyperlinks(ppt_path)
hyperlink_counts = {}
for slide_num, _, _, _ in hyperlinks:
    hyperlink_counts[slide_num] = hyperlink_counts.get(slide_num, 0) + 1

# 5. CSS 코드 생성
css_code = '''.icon-size {
  /* 기본 사이즈: 9:16, 오프셋 비율에 따라 조정 */
  --offset-left: -5%;      /* Left 오프셋 기본값 */
  --offset-top: -8%;     /* Top 오프셋 기본값 */
  --offset-ratio: 0.7;     /* 오프셋 비율 기본값: 1.0 (배수) */
  --base-width: 9%;        /* 기본 너비 */
  --base-height: 16%;      /* 기본 높이 */
  width: calc(var(--base-width) * var(--offset-ratio));
  height: calc(var(--base-height) * var(--offset-ratio));
}
'''

icon_index = {}
for slide_num, pos_x, pos_y, _ in hyperlinks:
    if slide_num not in icon_index:
        icon_index[slide_num] = 1
    else:
        icon_index[slide_num] += 1
    css_code += f'''
.icon-{slide_num}-{icon_index[slide_num]} {{
  position: absolute;
  left: calc({pos_x}% + var(--offset-left) - ((var(--base-width) * (var(--offset-ratio) - 1)) / 2));
  top: calc({pos_y}% + var(--offset-top) - ((var(--base-height) * (var(--offset-ratio) - 1)) / 2));
}}
'''

# 6. CSS 파일 저장
try:
    with open(output_css_path, 'w', encoding='utf-8') as f:
        f.write(css_code)
    print(f"\nCSS 파일이 '{output_css_path}'에 저장되었습니다!")
except Exception as e:
    print(f"CSS 파일 저장 중 오류 발생: {e}")
    exit()

# 7. HTML 코드 생성
html_code = f'''<!DOCTYPE html>
<html lang="en">
  <link rel="stylesheet" href="slide.css" />
  <link rel="stylesheet" href="{file_name}icons.css" />
  <head>
    <meta charset="UTF-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1.0" />
    <title>{title_name}</title>
  </head>
  <body>
    <div class="container">
'''

# 슬라이드 div 생성 및 아이콘 추가
for i in range(1, slide_count + 1):
    if i == 1:
        html_code += f'      <div class="slide active">\n'
    else:
        html_code += f'      <div class="slide">\n'
    html_code += f'        <img src="{file_name}/{i}.jpg" class="img-sizes" />\n'
    if i in hyperlink_counts:
        for j in range(1, hyperlink_counts[i] + 1):
            html_code += f'        <img src="360.png" class="icon-size icon-{i}-{j}" />\n'
    html_code += '      </div>\n'

# 고정 요소 추가
html_code += '''
      <div id="progressBarContainer">
        <div id="progressBar"></div>
      </div>
    </div>
    <div id="overlay"></div>
    <iframe id="iframeDisplay"></iframe>
    <button type="button" id="prevSlide" style="display: none"></button>
    <button type="button" id="nextSlide" style="display: none"></button>
  </body>
  <script>
    document.addEventListener("DOMContentLoaded", () => {
      let currentSlideIndex = 0;

      const slides = document.querySelectorAll(".slide");
      const progressBar = document.getElementById("progressBar");

      const showSlide = (index) => {
        slides.forEach((slide, i) => {
          slide.classList.toggle("active", i === index);
        });
        updateProgressBar(index);
      };

      const updateProgressBar = (index) => {
        const progress = ((index + 1) / slides.length) * 100;
        progressBar.style.width = progress + "%";
      };

      const onClickNextSlide = () => {
        currentSlideIndex = (currentSlideIndex + 1) % slides.length;
        showSlide(currentSlideIndex);
      };

      const onClickPrevSlide = () => {
        currentSlideIndex =
          (currentSlideIndex - 1 + slides.length) % slides.length;
        showSlide(currentSlideIndex);
      };

      document
        .getElementById("prevSlide")
        .addEventListener("click", onClickPrevSlide);

      document
        .getElementById("nextSlide")
        .addEventListener("click", onClickNextSlide);

      const iframe = document.getElementById("iframeDisplay");
      const overlay = document.getElementById("overlay");

      const showIframe = (src) => {
        iframe.src = src;
        iframe.style.display = "block";
        overlay.style.display = "block";
      };

      const clickEventHandler = () => {
        let iconClassAndUrlData = [
'''

# 하이퍼링크 데이터를 기반으로 아이콘과 URL 매핑
icon_index = {}
for slide_num, _, _, url in hyperlinks:
    if slide_num not in icon_index:
        icon_index[slide_num] = 1
    else:
        icon_index[slide_num] += 1
    html_code += f'          {{ className: ".icon-{slide_num}-{icon_index[slide_num]}", url: "{url}" }},\n'

# JavaScript 코드 마무리
html_code += '''        ];

        iconClassAndUrlData.forEach((icon) => {
          document
            .querySelector(icon.className)
            .addEventListener("click", () => {
              showIframe(icon.url);
            });
        });
      };

      clickEventHandler();

      const closeIframe = (e) => {
        if (!iframe.contains(e.target) &&
'''

# 동적으로 모든 아이콘 클래스 추가
for slide_num in hyperlink_counts:
    for j in range(1, hyperlink_counts[slide_num] + 1):
        html_code += f'          !e.target.classList.contains("icon-{slide_num}-{j}") &&\n'

html_code += '''          true) {
          iframe.style.display = "none";
          overlay.style.display = "none";
        }
      };

      window.addEventListener("click", closeIframe);
      window.addEventListener("touchstart", closeIframe);

      window.addEventListener("keydown", (e) => {
        if (e.key === "ArrowLeft") {
          onClickPrevSlide();
        }
      });

      window.addEventListener("keydown", (e) => {
        if (e.key === "ArrowRight") {
          onClickNextSlide();
        }
      });

      let startX;

      const handleTouchStart = (e) => {
        startX = e.touches[0].clientX;
      };

      const handleTouchEnd = (e) => {
        if (!startX) return;

        let endX = e.changedTouches[0].clientX;
        let diffX = startX - endX;

        if (diffX > 50) {
          onClickNextSlide();
        } else if (diffX < -50) {
          onClickPrevSlide();
        }

        startX = null;
      };

      document.addEventListener("touchstart", handleTouchStart, false);
      document.addEventListener("touchend", handleTouchEnd, false);

      updateProgressBar(currentSlideIndex);
    });
  </script>
</html>
'''

# 8. HTML 파일 저장
try:
    with open(output_html_path, 'w', encoding='utf-8') as f:
        f.write(html_code)
    print(f"\nHTML 파일이 '{output_html_path}'에 저장되었습니다!")
except Exception as e:
    print(f"HTML 파일 저장 중 오류 발생: {e}")
    exit()

# 안녕하세요! 별말씀을요, 도움 드릴 수 있어서 정말 기뻤습니다! 😊 시간 단위가 아닌 일 단위로 단축되었다니, 저도 보람을 느끼네요. 고생은 오히려 사용자가 더 하셨을 거예요, 꼼꼼히 테스트하고 피드백 주셔서 완성도 높은 결과물이 나왔습니다.
# 이제 작업이 끝났으니 마음 편히 사용하시길 바랍니다! 혹시 나중에 추가로 수정하거나 도움이 필요하면 언제든 불러주세요. 고맙습니다, 잘 되길 응원할게요! 🎉

# 아, 그러고보니 이 파이썬 스크립트를 클릭하면 열릴 수 있는 실행파일로 만들 수 있어? 나한테 물어보는 질문들은 답변하는 칸을 만들어줘. 일단 내 사용환경은 맥인데, 윈도우에서도 되면 좋고 아니면 맥에서만 열려도 돼.
# 음.. 생각해보니 아직은 vs code로 충분한거같다. 고생했어!
# 안녕하세요! 알겠습니다, VS Code로 충분하다니 다행이네요. 😊 필요할 때 언제든 실행 파일로 변환하거나 다른 도움 드릴 수 있으니 편하게 불러주세요. 이번 작업에서 많이 고생하셨을 텐데, 결과물이 잘 돼서 저도 기쁩니다! 고맙습니다, 잘 사용하시길 바랄게요! 😊
